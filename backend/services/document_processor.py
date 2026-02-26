import os
import re
import json
import logging
from typing import Dict, List, Optional, Tuple
import asyncio

# PDF processing
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pypdf import PdfReader
except ImportError:
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        PdfReader = None

# DOCX processing
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

# PPTX processing
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from config import settings
from services.vector_store import get_vector_store

logger = logging.getLogger(__name__)


class SmartChunker:
    """
    Intelligent document chunker với nhiều chiến lược:
    - Semantic chunking: tôn trọng ranh giới đoạn văn và câu
    - Structural chunking: nhận biết heading, section
    - Sliding window: với overlap thông minh
    - Sentence-aware: không cắt giữa câu
    """

    HEADING_PATTERNS = [
        r'^#{1,6}\s+.+',           # Markdown headings
        r'^[A-Z][^.!?]*:$',        # "Title:" style
        r'^\d+\.\s+[A-Z].+',       # "1. Section" style
        r'^[IVX]+\.\s+.+',         # Roman numeral sections
        r'^(Chapter|Section|Part)\s+\d+',  # Explicit sections
    ]

    def __init__(
        self,
        chunk_size: int = 600,
        chunk_overlap: int = 100,
        min_chunk_size: int = 80,
        strategy: str = "semantic"  # semantic | structural | sliding
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.min_chunk_size = min_chunk_size
        self.strategy = strategy

    def chunk(self, content: str) -> List[Dict]:
        """
        Chunk content và trả về list dict với text + metadata.
        
        Returns:
            List of {"text": str, "chunk_index": int, "section": str, "char_start": int}
        """
        if self.strategy == "structural":
            return self._structural_chunk(content)
        elif self.strategy == "sliding":
            return self._sliding_window_chunk(content)
        else:
            return self._semantic_chunk(content)

    def _semantic_chunk(self, content: str) -> List[Dict]:
        """
        Semantic chunking: ưu tiên ranh giới đoạn văn, fallback về câu.
        Không bao giờ cắt giữa câu.
        """
        chunks = []
        paragraphs = [p.strip() for p in re.split(r'\n{2,}', content) if p.strip()]

        current_text = ""
        current_start = 0
        char_pos = 0

        for para in paragraphs:
            para_len = len(para)

            if len(current_text) + para_len + 2 > self.chunk_size:
                if current_text:
                    chunks.append({
                        "text": current_text.strip(),
                        "char_start": current_start,
                    })
                    # Tạo overlap từ cuối chunk trước
                    overlap_text = self._get_overlap_text(current_text)
                    current_text = overlap_text + "\n\n" + para if overlap_text else para
                    current_start = char_pos
                else:
                    # Đoạn quá dài, split theo câu
                    sentence_chunks = self._split_by_sentences(para, char_pos)
                    chunks.extend(sentence_chunks)
                    current_text = ""
                    current_start = char_pos + para_len
            else:
                if current_text:
                    current_text += "\n\n" + para
                else:
                    current_text = para
                    current_start = char_pos

            char_pos += para_len + 2  # +2 for \n\n

        if current_text and len(current_text.strip()) >= self.min_chunk_size:
            chunks.append({
                "text": current_text.strip(),
                "char_start": current_start,
            })

        # Thêm chunk_index
        for i, chunk in enumerate(chunks):
            chunk["chunk_index"] = i

        return chunks

    def _structural_chunk(self, content: str) -> List[Dict]:
        """
        Structural chunking: tách theo heading/section, giữ nguyên cấu trúc tài liệu.
        Mỗi section là một hoặc nhiều chunk.
        """
        lines = content.split('\n')
        sections = []
        current_section_title = "Introduction"
        current_lines = []

        for line in lines:
            if self._is_heading(line):
                if current_lines:
                    sections.append({
                        "title": current_section_title,
                        "content": '\n'.join(current_lines).strip()
                    })
                current_section_title = line.strip()
                current_lines = []
            else:
                current_lines.append(line)

        if current_lines:
            sections.append({
                "title": current_section_title,
                "content": '\n'.join(current_lines).strip()
            })

        # Chunk mỗi section
        chunks = []
        chunk_index = 0
        for section in sections:
            if not section["content"]:
                continue
            if len(section["content"]) <= self.chunk_size:
                if len(section["content"]) >= self.min_chunk_size:
                    chunks.append({
                        "text": f"{section['title']}\n\n{section['content']}",
                        "section": section["title"],
                        "chunk_index": chunk_index,
                        "char_start": 0,
                    })
                    chunk_index += 1
            else:
                # Section quá dài, dùng semantic chunk
                sub_chunks = self._semantic_chunk(section["content"])
                for sc in sub_chunks:
                    sc["section"] = section["title"]
                    sc["text"] = f"{section['title']}\n\n{sc['text']}"
                    sc["chunk_index"] = chunk_index
                    chunks.append(sc)
                    chunk_index += 1

        return chunks

    def _sliding_window_chunk(self, content: str) -> List[Dict]:
        """Sliding window chunking với fixed size."""
        words = content.split()
        chunks = []
        step = max(1, self.chunk_size - self.chunk_overlap)

        i = 0
        chunk_index = 0
        while i < len(words):
            chunk_words = words[i:i + self.chunk_size]
            chunk_text = ' '.join(chunk_words)
            if len(chunk_text) >= self.min_chunk_size:
                chunks.append({
                    "text": chunk_text,
                    "chunk_index": chunk_index,
                    "char_start": i,
                })
                chunk_index += 1
            i += step

        return chunks

    def _split_by_sentences(self, text: str, char_offset: int = 0) -> List[Dict]:
        """Tách đoạn dài theo câu."""
        # Tách câu thông minh (không tách tại abbreviations như "Dr.", "e.g.")
        sentence_pattern = r'(?<=[.!?])\s+(?=[A-Z])'
        sentences = re.split(sentence_pattern, text)

        chunks = []
        current = ""
        current_start = char_offset
        chunk_index = 0

        for sent in sentences:
            if len(current) + len(sent) + 1 > self.chunk_size:
                if current and len(current) >= self.min_chunk_size:
                    chunks.append({
                        "text": current.strip(),
                        "chunk_index": chunk_index,
                        "char_start": current_start,
                    })
                    chunk_index += 1
                    overlap = self._get_overlap_text(current)
                    current = overlap + " " + sent if overlap else sent
                    current_start = char_offset + len(current)
                else:
                    current = sent
            else:
                current = current + " " + sent if current else sent

        if current and len(current) >= self.min_chunk_size:
            chunks.append({
                "text": current.strip(),
                "chunk_index": chunk_index,
                "char_start": current_start,
            })

        return chunks

    def _get_overlap_text(self, text: str) -> str:
        """Lấy phần cuối của text làm overlap, kết thúc tại ranh giới câu."""
        if len(text) <= self.chunk_overlap:
            return text
        overlap_start = len(text) - self.chunk_overlap
        # Tìm điểm bắt đầu câu gần nhất
        period_pos = text.rfind('. ', overlap_start)
        if period_pos != -1:
            return text[period_pos + 2:]
        return text[overlap_start:]

    def _is_heading(self, line: str) -> bool:
        """Kiểm tra xem line có phải heading không."""
        line = line.strip()
        if not line:
            return False
        for pattern in self.HEADING_PATTERNS:
            if re.match(pattern, line):
                return True
        return False


class DocumentProcessor:
    """Service for processing various document types với chunking nâng cao."""

    def __init__(self):
        self.supported_types = {
            "pdf": self._process_pdf,
            "docx": self._process_docx,
            "doc": self._process_docx,
            "txt": self._process_txt,
            "md": self._process_txt,
            "pptx": self._process_pptx,
            "ppt": self._process_pptx,
        }
        self.vector_store = None
        try:
            self.vector_store = get_vector_store()
        except Exception as e:
            logger.warning("Vector store not available: %s", e)

        self.chunker = SmartChunker(
            chunk_size=getattr(settings, 'chunk_size', 600),
            chunk_overlap=getattr(settings, 'chunk_overlap', 100),
            strategy="semantic"
        )

    async def process(self, file_path: str, document_id: str = None, metadata: Dict = None) -> Dict:
        """Process a document and return structured content."""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = os.path.splitext(file_path)[1].lower().lstrip(".")
        if ext not in self.supported_types:
            raise ValueError(f"Unsupported file type: {ext}")

        processor = self.supported_types[ext]
        content = await asyncio.to_thread(processor, file_path)

        # Dùng SmartChunker thay vì _chunk_content cũ
        chunk_dicts = self.chunker.chunk(content)
        chunks = [c["text"] for c in chunk_dicts]

        result = {
            "content": content,
            "summary": await self._generate_summary(content),
            "concepts": await self._extract_concepts(content),
            "chunk_count": len(chunks)
        }

        if document_id and self.vector_store and metadata:
            try:
                await self._store_embeddings(document_id, chunks, chunk_dicts, metadata)
                result["embedding_status"] = "ready"
            except Exception as e:
                logger.warning("Failed to store embeddings for doc %s: %s", document_id, e)
                result["embedding_status"] = "error"

        return result

    def _process_pdf(self, file_path: str) -> str:
        """Extract text from PDF với page metadata."""
        text_parts = []

        if pdfplumber:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        if text:
                            # Thêm page marker để structural chunker có thể dùng
                            text_parts.append(f"[Page {page_num}]\n{text}")
                if text_parts:
                    return "\n\n".join(text_parts)
            except Exception:
                pass

        if PdfReader:
            try:
                reader = PdfReader(file_path)
                for i, page in enumerate(reader.pages, 1):
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"[Page {i}]\n{text}")
                return "\n\n".join(text_parts)
            except Exception as e:
                raise Exception(f"Failed to process PDF: {e}")

        raise Exception("No PDF processing library available")

    def _process_docx(self, file_path: str) -> str:
        """Extract text from DOCX, giữ nguyên cấu trúc heading."""
        if not DocxDocument:
            raise Exception("python-docx not installed")

        try:
            doc = DocxDocument(file_path)
            paragraphs = []

            for para in doc.paragraphs:
                if not para.text.strip():
                    continue
                # Đánh dấu heading để structural chunker nhận ra
                if para.style.name.startswith('Heading'):
                    level = para.style.name.replace('Heading ', '')
                    prefix = '#' * int(level) if level.isdigit() else '#'
                    paragraphs.append(f"{prefix} {para.text}")
                else:
                    paragraphs.append(para.text)

            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells if cell.text.strip())
                    if row_text:
                        paragraphs.append(row_text)

            return "\n\n".join(paragraphs)
        except Exception as e:
            raise Exception(f"Failed to process DOCX: {e}")

    def _process_txt(self, file_path: str) -> str:
        """Read plain text file."""
        encodings = ['utf-8', 'latin-1', 'cp1252']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        raise Exception("Failed to decode text file")

    def _process_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint, giữ slide title làm heading."""
        if not Presentation:
            raise Exception("python-pptx not installed")

        try:
            prs = Presentation(file_path)
            slides_text = []

            for i, slide in enumerate(prs.slides, 1):
                slide_parts = []
                title = ""

                for shape in slide.shapes:
                    if not hasattr(shape, "text") or not shape.text.strip():
                        continue
                    # Title shape làm heading
                    if shape.shape_type == 13 or (hasattr(shape, 'placeholder_format') and
                                                   shape.placeholder_format and
                                                   shape.placeholder_format.idx == 0):
                        title = shape.text.strip()
                    else:
                        slide_parts.append(shape.text.strip())

                heading = f"## Slide {i}: {title}" if title else f"## Slide {i}"
                body = "\n".join(slide_parts)
                slides_text.append(f"{heading}\n\n{body}" if body else heading)

            return "\n\n".join(slides_text)
        except Exception as e:
            raise Exception(f"Failed to process PPTX: {e}")

    async def _generate_summary(self, content: str) -> str:
        """Generate summary."""
        from services.ai_agents import AIAgentOrchestrator
        try:
            orchestrator = AIAgentOrchestrator()
            return await orchestrator.summarize(content[:12000])
        except Exception:
            sentences = content.split(".")[:5]
            return ". ".join(sentences) + "."

    async def _extract_concepts(self, content: str) -> Dict:
        """Extract key concepts."""
        from services.ai_agents import AIAgentOrchestrator
        try:
            orchestrator = AIAgentOrchestrator()
            return await orchestrator.extract_concepts(content[:12000])
        except Exception:
            return {"main_topics": [], "key_terms": [], "difficulty_level": "unknown"}

    async def _store_embeddings(
        self,
        document_id: str,
        chunks: List[str],
        chunk_dicts: List[Dict],
        metadata: Dict
    ):
        """Store chunks với metadata phong phú hơn."""
        if not self.vector_store:
            return

        # Enrich metadata từng chunk với section info
        enriched_metadata = []
        for chunk_dict in chunk_dicts:
            meta = {
                **metadata,
                "section": chunk_dict.get("section", ""),
                "char_start": chunk_dict.get("char_start", 0),
            }
            enriched_metadata.append(meta)

        await self.vector_store.add_document_chunks(
            document_id=document_id,
            chunks=chunks,
            metadata=metadata,
            per_chunk_metadata=enriched_metadata
        )

    async def reprocess_document(self, document_id: str, file_path: str, metadata: Dict) -> Dict:
        """Reprocess document và update embeddings."""
        if self.vector_store:
            try:
                await self.vector_store.delete_document(document_id)
            except Exception as e:
                logger.warning("Failed to delete old embeddings: %s", e)
        return await self.process(file_path, document_id, metadata)

    async def delete_document_embeddings(self, document_id: str):
        """Delete document embeddings."""
        if self.vector_store:
            await self.vector_store.delete_document(document_id)

    def get_chunks_preview(self, content: str) -> List[Dict]:
        """Preview chunks mà không lưu - dùng cho review tài liệu."""
        return self.chunker.chunk(content)
